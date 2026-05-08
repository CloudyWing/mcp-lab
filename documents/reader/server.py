#!/usr/bin/env python3
"""
mcp-reader — Document & image reader MCP server.
"""

import base64
import chardet
import json
import os
import re
from collections import deque
from pathlib import Path
from typing import Any

from mcp.server.fastmcp import FastMCP

MOUNT_ROOT = Path("/host_root").resolve()
RAW_HOST_MOUNT_ROOT = os.environ.get("DOCUMENT_HOST_MOUNT_ROOT", "/").strip()


def _get_int_env(name: str, default: int) -> int:
    try:
        return int(os.environ.get(name, default))
    except (TypeError, ValueError):
        return default


def _get_float_env(name: str, default: float) -> float:
    try:
        return float(os.environ.get(name, default))
    except (TypeError, ValueError):
        return default


MAX_IMAGES = _get_int_env("MAX_IMAGES", 20)
MAX_TEXT_MB = _get_float_env("MAX_TEXT_MB", 10)
DEFAULT_LIST_LIMIT = _get_int_env("DOCUMENT_LIST_DEFAULT_LIMIT", 100)
MAX_LIST_LIMIT = _get_int_env("DOCUMENT_LIST_MAX_LIMIT", 1000)
DEFAULT_LIST_DEPTH = _get_int_env("DOCUMENT_LIST_DEFAULT_DEPTH", 2)
MAX_LIST_DEPTH = _get_int_env("DOCUMENT_LIST_MAX_DEPTH", 10)
DEFAULT_SEARCH_LIMIT = _get_int_env("DOCUMENT_SEARCH_DEFAULT_LIMIT", 20)
MAX_SEARCH_LIMIT = _get_int_env("DOCUMENT_SEARCH_MAX_LIMIT", 200)
DEFAULT_SEARCH_SCAN_LIMIT = _get_int_env("DOCUMENT_SEARCH_DEFAULT_SCAN_LIMIT", 200)
MAX_SEARCH_SCAN_LIMIT = _get_int_env("DOCUMENT_SEARCH_MAX_SCAN_LIMIT", 2000)
DEFAULT_XLSX_MAX_ROWS = _get_int_env("DOCUMENT_XLSX_MAX_ROWS", 1000)
MAX_XLSX_MAX_ROWS = _get_int_env("DOCUMENT_XLSX_MAX_ROWS_LIMIT", 10000)
DEFAULT_XLSX_MAX_CELLS = _get_int_env("DOCUMENT_XLSX_MAX_CELLS", 20000)
MAX_XLSX_MAX_CELLS = _get_int_env("DOCUMENT_XLSX_MAX_CELLS_LIMIT", 200000)

DOCUMENT_EXTS = {".pdf", ".pptx", ".docx", ".xlsx"}
IMAGE_EXTS    = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".svg"}
TEXT_EXTS     = {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".log", ".ini", ".toml", ".html"}
SUPPORTED_EXTENSIONS = DOCUMENT_EXTS | IMAGE_EXTS | TEXT_EXTS

mcp = FastMCP(
    "mcp-reader",
    host="0.0.0.0",
    port=int(os.environ.get("PORT", 8080)),
)


def _normalize_mount_root(raw_root: str) -> str:
    normalized = raw_root.strip().strip('"').strip("'").replace("\\", "/")
    if not normalized:
        return "/"
    if not normalized.startswith("/"):
        normalized = f"/{normalized}"
    normalized = normalized.rstrip("/")
    return normalized or "/"


HOST_MOUNT_ROOT = _normalize_mount_root(RAW_HOST_MOUNT_ROOT)


def _clamp_int(value: int, default: int, minimum: int, maximum: int) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        return default

    if parsed < minimum:
        return default
    return min(parsed, maximum)


def _map_host_path(source_path: str) -> Path | None:
    if source_path == HOST_MOUNT_ROOT:
        return MOUNT_ROOT

    prefix = f"{HOST_MOUNT_ROOT}/"
    if source_path.startswith(prefix):
        tail = source_path.removeprefix(prefix).strip("/")
        return MOUNT_ROOT / tail if tail else MOUNT_ROOT

    return None


def _resolve_path(path: str = ".") -> Path:
    normalized = path.strip().strip('"').strip("'")
    if not normalized or normalized == ".":
        target = MOUNT_ROOT
    elif normalized.startswith("/host_root/") or normalized == "/host_root":
        target = Path(normalized).resolve()
    else:
        windows_path = re.match(r"^(?P<drive>[A-Za-z]):[\\/]*(?P<tail>.*)$", normalized)
        mapped_target: Path | None = None
        if windows_path:
            drive = windows_path.group("drive").lower()
            tail = windows_path.group("tail").replace("\\", "/").strip("/")
            host_style_path = f"/mnt/{drive}"
            if tail:
                host_style_path = f"{host_style_path}/{tail}"
            mapped_target = _map_host_path(host_style_path)
        else:
            normalized = normalized.replace("\\", "/")
            if normalized.startswith("/"):
                mapped_target = _map_host_path(normalized)
            else:
                mapped_target = (MOUNT_ROOT / normalized).resolve()

        if mapped_target is None:
            raise ValueError(f"Path must stay under DOCUMENT_HOST_MOUNT_ROOT: {HOST_MOUNT_ROOT}")
        target = mapped_target.resolve()

    try:
        target.relative_to(MOUNT_ROOT)
    except ValueError as exc:
        raise ValueError(f"Path must stay under DOCUMENT_HOST_MOUNT_ROOT: {HOST_MOUNT_ROOT}") from exc
    return target


def _relative_path_for(file_path: Path) -> str:
    try:
        relative_path = file_path.relative_to(MOUNT_ROOT).as_posix()
        if not relative_path or relative_path == ".":
            return HOST_MOUNT_ROOT
        if HOST_MOUNT_ROOT == "/":
            return f"/{relative_path}"
        return f"{HOST_MOUNT_ROOT}/{relative_path}"
    except Exception:
        return str(file_path.name)


def _candidate_host_path(file_path: Path) -> str:
    try:
        return _relative_path_for(file_path)
    except Exception:
        return str(file_path)


def _file_category(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext in DOCUMENT_EXTS:
        return "document"
    if ext in IMAGE_EXTS:
        return "image"
    return "text"


def _document_file_info(file_path: Path) -> dict:
    return {
        "name":     file_path.name,
        "path":     _relative_path_for(file_path),
        "type":     file_path.suffix.lstrip(".").lower(),
        "category": _file_category(file_path),
        "size_kb":  round(file_path.stat().st_size / 1024, 1),
    }


def _list_supported_files(target: Path, recursive: bool, max_depth: int, limit: int) -> tuple[list[dict], bool, bool]:
    if target.is_file():
        if target.suffix.lower() not in SUPPORTED_EXTENSIONS:
            return [], False, False
        return [_document_file_info(target)], False, False

    if not target.is_dir():
        return [], False, False

    files: list[dict] = []
    depth_limited = False
    queue = deque([(target, 0)])

    while queue:
        current, depth = queue.popleft()

        try:
            entries = sorted(current.iterdir(), key=lambda item: item.name.lower())
        except OSError:
            continue

        for entry in entries:
            if entry.is_file() and entry.suffix.lower() in SUPPORTED_EXTENSIONS:
                if len(files) >= limit:
                    return files, True, depth_limited
                files.append(_document_file_info(entry))
                continue

            if recursive and entry.is_dir() and not entry.is_symlink():
                if depth < max_depth:
                    queue.append((entry, depth + 1))
                else:
                    depth_limited = True

    return files, False, depth_limited


def _find_path_candidates(raw_path: str, limit: int = 10) -> list[str]:
    try:
        requested = _resolve_path(raw_path)
    except ValueError:
        return []

    if requested.exists():
        return [_candidate_host_path(requested)]

    parent = requested.parent
    name_prefix = requested.name.lower()
    if not parent.exists():
        return []

    candidates: list[str] = []
    for file_path in sorted(parent.iterdir(), key=lambda item: item.name.lower()):
        if not file_path.is_file():
            continue
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if name_prefix and not file_path.name.lower().startswith(name_prefix):
            continue
        candidates.append(_candidate_host_path(file_path))
        if len(candidates) >= limit:
            break
    return candidates


def _img_to_b64(img_bytes: bytes, fmt: str = "png") -> str:
    return f"data:image/{fmt};base64," + base64.b64encode(img_bytes).decode()


def _parse_pdf(path: Path, include_images: bool, start_page: int, end_page: int) -> dict:
    import fitz

    document = fitz.open(str(path))
    total = len(document)
    p_start = max(0, start_page - 1)
    p_end   = (min(end_page, total) - 1) if end_page > 0 else (total - 1)
    pages, images_b64 = [], []
    for i in range(p_start, p_end + 1):
        page     = document[i]
        page_num = i + 1
        text     = page.get_text("text").strip()
        if include_images and len(images_b64) < MAX_IMAGES:
            matrix = fitz.Matrix(1.5, 1.5)
            pixmap = page.get_pixmap(matrix=matrix)
            images_b64.append({"page": page_num, "image": _img_to_b64(pixmap.tobytes("png"))})
        pages.append({"page": page_num, "text": text})
    document.close()
    return {
        "type": "pdf",
        "total_pages": total,
        "returned_pages": f"{p_start + 1}-{p_end + 1}",
        "pages": pages,
        "images": images_b64,
    }


def _parse_pptx(path: Path, include_images: bool, start_page: int, end_page: int) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs   = Presentation(str(path))
    total = len(prs.slides)
    s_start = max(0, start_page - 1)
    s_end   = (min(end_page, total) - 1) if end_page > 0 else (total - 1)
    slides, images_b64 = [], []
    for i in range(s_start, s_end + 1):
        slide     = prs.slides[i]
        slide_num = i + 1
        texts, slide_images = [], []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            if include_images and len(images_b64) < MAX_IMAGES and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    img_fmt   = shape.image.ext.lower().lstrip(".")
                    if img_fmt == "jpg":
                        img_fmt = "jpeg"
                    encoded = _img_to_b64(img_bytes, img_fmt)
                    slide_images.append(encoded)
                    images_b64.append({"slide": slide_num, "image": encoded})
                except Exception:
                    pass
        title = slide.shapes.title.text if slide.shapes.title else ""
        slides.append({"slide": slide_num, "title": title, "texts": texts, "images": slide_images})
    return {
        "type": "pptx",
        "total_slides": total,
        "returned_slides": f"{s_start + 1}-{s_end + 1}",
        "slides": slides,
        "images": images_b64,
    }


def _parse_docx(path: Path, include_images: bool) -> dict:
    import docx

    doc        = docx.Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    images_b64 = []
    if include_images:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype and len(images_b64) < MAX_IMAGES:
                try:
                    img_bytes = rel.target_part.blob
                    img_fmt   = rel.target_part.content_type.split("/")[-1]
                    if img_fmt == "jpg":
                        img_fmt = "jpeg"
                    images_b64.append({"index": len(images_b64) + 1, "image": _img_to_b64(img_bytes, img_fmt)})
                except Exception:
                    pass
    tables = []
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append(rows)
    return {"type": "docx", "paragraphs": paragraphs, "tables": tables, "images": images_b64}


def _parse_xlsx(path: Path, max_rows: int, max_cells: int) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True)
    total_sheet_count = len(wb.sheetnames)
    sheets = []
    total_rows = 0
    total_cells = 0
    truncated = False

    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if total_rows >= max_rows or total_cells + len(row) > max_cells:
                truncated = True
                break

            rows.append([str(c) if c is not None else "" for c in row])
            total_rows += 1
            total_cells += len(row)

        sheets.append({"sheet": name, "rows": rows})

        if truncated:
            break

    wb.close()
    return {
        "type": "xlsx",
        "sheet_count": total_sheet_count,
        "returned_sheet_count": len(sheets),
        "returned_rows": total_rows,
        "truncated": truncated,
        "sheets": sheets,
        "images": [],
    }


def _parse_image(path: Path) -> dict:
    from PIL import Image
    import io

    with Image.open(str(path)) as img:
        width, height = img.size
        fmt      = img.format or path.suffix.lstrip(".").upper()
        mime_fmt = fmt.lower()
        if mime_fmt in ("jpg",):
            mime_fmt = "jpeg"

        buf = io.BytesIO()
        save_fmt = "PNG" if fmt in ("", "ICO", "ICNS") else fmt
        try:
            img.save(buf, format=save_fmt)
        except Exception:
            img.convert("RGB").save(buf, format="PNG")
            mime_fmt = "png"
        image_b64 = _img_to_b64(buf.getvalue(), mime_fmt)

    result: dict[str, Any] = {
        "type":       "image",
        "format":     fmt,
        "width":      width,
        "height":     height,
        "size_bytes": path.stat().st_size,
        "image":      image_b64,
        "images":     [],
    }

    try:
        import pytesseract
        ocr_text = pytesseract.image_to_string(str(path), lang="chi_tra+chi_sim+eng").strip()
        if ocr_text:
            result["ocr_text"] = ocr_text
    except Exception:
        pass

    return result


def _parse_text(path: Path) -> dict:
    size_bytes = path.stat().st_size
    max_bytes  = int(MAX_TEXT_MB * 1024 * 1024) if MAX_TEXT_MB > 0 else 0
    truncated  = max_bytes > 0 and size_bytes > max_bytes

    raw = path.read_bytes() if max_bytes <= 0 else path.read_bytes()[:max_bytes]
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"
    try:
        content = raw.decode(encoding, errors="replace")
    except Exception:
        content = raw.decode("utf-8", errors="replace")

    return {
        "type":       "text",
        "encoding":   encoding,
        "line_count": len(content.splitlines()),
        "size_bytes": size_bytes,
        "truncated":  truncated,
        "content":    content,
        "images":     [],
    }


@mcp.tool()
def get_document_root() -> dict:
    """回傳目前掛載根目錄資訊、支援格式，與路徑傳入方式提示。"""
    return {
        "document_root": str(MOUNT_ROOT),
        "host_mount_root": HOST_MOUNT_ROOT,
        "path_usage": {
            "json_encoding": "utf-8",
            "accepts": [
                "完整 Windows 路徑，例如 C:\\Workspace\\WSL\\Reader\\sample.txt",
                "完整 WSL 路徑，例如 /home/user/documents/sample.txt",
                "掛載根目錄下相對路徑，例如 sample.txt",
            ],
            "partial_path_strategy": "若只有部分檔名，先呼叫 resolve_document_path 取得候選完整路徑，再呼叫 read_document。",
            "requires_under_mount_root": True,
        },
        "supported_extensions": sorted(SUPPORTED_EXTENSIONS),
        "categories": {
            "documents": sorted(DOCUMENT_EXTS),
            "images":    sorted(IMAGE_EXTS),
            "text":      sorted(TEXT_EXTS),
        },
    }


@mcp.tool()
def resolve_document_path(path: str, limit: int = 10) -> dict:
    """解析指定路徑是否可讀；若檔案不存在，回傳同目錄下符合前綴的候選檔案。"""
    try:
        resolved = _resolve_path(path)
    except ValueError as exc:
        return {
            "input_path": path,
            "host_mount_root": HOST_MOUNT_ROOT,
            "exists": False,
            "error": str(exc),
            "candidates": [],
        }

    candidates = [] if resolved.exists() else _find_path_candidates(path, limit=limit)
    return {
        "input_path": path,
        "resolved_path": str(resolved),
        "host_mount_root": HOST_MOUNT_ROOT,
        "exists": resolved.exists(),
        "is_file": resolved.is_file() if resolved.exists() else False,
        "candidates": candidates,
    }


@mcp.tool()
def list_documents(
    subpath: str = ".",
    recursive: bool = False,
    limit: int = DEFAULT_LIST_LIMIT,
    max_depth: int = DEFAULT_LIST_DEPTH,
) -> dict:
    """列出掛載根目錄下可讀取的檔案清單，預設不遞迴且有回傳上限。"""

    safe_limit = _clamp_int(limit, DEFAULT_LIST_LIMIT, 1, MAX_LIST_LIMIT)
    safe_depth = _clamp_int(max_depth, DEFAULT_LIST_DEPTH, 0, MAX_LIST_DEPTH)
    response = {
        "subpath": subpath,
        "host_mount_root": HOST_MOUNT_ROOT,
        "recursive": recursive,
        "max_depth": safe_depth if recursive else 0,
        "limit": safe_limit,
        "exists": False,
        "returned_count": 0,
        "depth_limited": False,
        "truncated": False,
        "files": [],
    }

    try:
        target = _resolve_path(subpath)
    except ValueError as exc:
        response["error"] = str(exc)
        return response

    if not target.exists():
        return response

    files, truncated, depth_limited = _list_supported_files(target, recursive, safe_depth, safe_limit)
    response["exists"] = True
    response["returned_count"] = len(files)
    response["depth_limited"] = depth_limited
    response["truncated"] = truncated or depth_limited
    response["files"] = files
    return response


@mcp.tool()
def read_document(
    path: str,
    include_images: bool = True,
    start_page: int = 1,
    end_page: int = 0,
    xlsx_max_rows: int = DEFAULT_XLSX_MAX_ROWS,
    xlsx_max_cells: int = DEFAULT_XLSX_MAX_CELLS,
) -> str:
    """讀取文件、圖片或文字檔，path 可傳完整 Windows/WSL 路徑或掛載根下相對路徑。
    start_page / end_page 僅適用於 PDF 和 PPTX（1-based，end_page=0 表示讀到末頁）。
    """

    try:
        file_path = _resolve_path(path)
    except ValueError as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)

    if not file_path.exists():
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return json.dumps({"error": f"Unsupported format: {ext}"}, ensure_ascii=False)

    try:
        if ext == ".pdf":
            result = _parse_pdf(file_path, include_images, start_page, end_page)
        elif ext == ".pptx":
            result = _parse_pptx(file_path, include_images, start_page, end_page)
        elif ext == ".docx":
            result = _parse_docx(file_path, include_images)
        elif ext == ".xlsx":
            safe_rows = _clamp_int(xlsx_max_rows, DEFAULT_XLSX_MAX_ROWS, 1, MAX_XLSX_MAX_ROWS)
            safe_cells = _clamp_int(xlsx_max_cells, DEFAULT_XLSX_MAX_CELLS, 1, MAX_XLSX_MAX_CELLS)
            result = _parse_xlsx(file_path, safe_rows, safe_cells)
        elif ext in IMAGE_EXTS:
            result = _parse_image(file_path)
        else:
            result = _parse_text(file_path)

        result["path"]        = _relative_path_for(file_path)
        result["image_count"] = len(result.get("images", []))
        return json.dumps(result, ensure_ascii=False)
    except Exception as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)


@mcp.tool()
def read_document_text_only(
    path: str,
    start_page: int = 1,
    end_page: int = 0,
) -> str:
    """只讀取文字（不擷取圖片，速度較快）。
    start_page / end_page 適用於 PDF 和 PPTX。
    """
    return read_document(path, include_images=False, start_page=start_page, end_page=end_page)


@mcp.tool()
def search_documents(
    keyword: str,
    subpath: str = ".",
    recursive: bool = True,
    limit: int = DEFAULT_SEARCH_LIMIT,
    max_files_scanned: int = DEFAULT_SEARCH_SCAN_LIMIT,
    max_depth: int = DEFAULT_LIST_DEPTH,
    include_image_ocr: bool = False,
) -> dict:
    """在指定路徑下搜尋包含關鍵字的檔案，預設限制掃描數量且不做圖片 OCR。"""
    safe_limit = _clamp_int(limit, DEFAULT_SEARCH_LIMIT, 1, MAX_SEARCH_LIMIT)
    safe_scan_limit = _clamp_int(max_files_scanned, DEFAULT_SEARCH_SCAN_LIMIT, 1, MAX_SEARCH_SCAN_LIMIT)
    safe_depth = _clamp_int(max_depth, DEFAULT_LIST_DEPTH, 0, MAX_LIST_DEPTH)
    response = {
        "keyword": keyword,
        "subpath": subpath,
        "recursive": recursive,
        "max_depth": safe_depth if recursive else 0,
        "limit": safe_limit,
        "max_files_scanned": safe_scan_limit,
        "scanned_count": 0,
        "returned_count": 0,
        "truncated": False,
        "results": [],
    }

    if not keyword.strip():
        response["error"] = "Keyword is required."
        return response

    listing = list_documents(
        subpath=subpath,
        recursive=recursive,
        limit=safe_scan_limit,
        max_depth=safe_depth,
    )

    if "error" in listing:
        response["error"] = listing["error"]
        return response

    results = []
    for doc in listing["files"]:
        if doc.get("category") == "image" and not include_image_ocr:
            continue

        response["scanned_count"] += 1

        try:
            content = json.loads(read_document_text_only(doc["path"]))
        except Exception:
            continue

        text_blob = json.dumps(content, ensure_ascii=False)
        if keyword.lower() in text_blob.lower():
            results.append({
                "path":     doc["path"],
                "type":     doc["type"],
                "category": doc.get("category", ""),
            })

            if len(results) >= safe_limit:
                response["truncated"] = True
                break

    response["returned_count"] = len(results)
    response["truncated"] = response["truncated"] or listing.get("truncated", False)
    response["results"] = results
    return response


if __name__ == "__main__":
    mcp.run(transport="streamable-http")
